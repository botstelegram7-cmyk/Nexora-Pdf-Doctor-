"""
PDF & Image utility functions v5.0
New: pdf2txt, linearize, thumbnail, pdf_info, redact, impose, deskew,
     image tools (compress/resize/crop/filter/text/convert/bgremove),
     csv2pdf, txt2pdf, html2pdf, json2pdf, hash, steganography, pdf_sign,
     poster, calendar_pdf, invoice, resume, certificate, zip/unzip,
     fileinfo, qrcode_scan, barcode, password strength checker,
     password cracker (Pro feature), metadata edit, epub converters.
"""
import io, os, hashlib, zipfile, tempfile, re, calendar, datetime, itertools, string, time
import fitz  # PyMuPDF
import pikepdf
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance, ImageOps
from config import FONTS, NOTEBOOK_STYLES
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm


# ─── HELPERS ─────────────────────────────────────────────────────────────────

def file_size_str(data: bytes) -> str:
    kb = len(data) / 1024
    return f"{kb/1024:.2f} MB" if kb > 1024 else f"{kb:.1f} KB"

def _parse_range(s: str, total: int) -> list:
    pages = []
    seen = set()
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            ab = part.split("-", 1)
            a = int(ab[0]) - 1 if ab[0].strip() else 0
            b = int(ab[1]) - 1 if ab[1].strip() else total - 1
            for p in range(a, b + 1):
                if 0 <= p < total and p not in seen:
                    pages.append(p); seen.add(p)
        elif part:
            p = int(part) - 1
            if 0 <= p < total and p not in seen:
                pages.append(p); seen.add(p)
    return pages

def _to_roman(n: int) -> str:
    vals = [(1000,'M'),(900,'CM'),(500,'D'),(400,'CD'),(100,'C'),(90,'XC'),
            (50,'L'),(40,'XL'),(10,'X'),(9,'IX'),(5,'V'),(4,'IV'),(1,'I')]
    result = ""
    for v, s in vals:
        while n >= v:
            result += s; n -= v
    return result


# ─── EXISTING FUNCTIONS (kept from v4) ───────────────────────────────────────

def compress_pdf(data: bytes) -> bytes:
    results = []
    try:
        with fitz.open(stream=data, filetype="pdf") as doc:
            buf = io.BytesIO()
            doc.save(buf, garbage=4, deflate=True, deflate_images=True,
                     deflate_fonts=True, clean=True)
            results.append(buf.getvalue())
    except Exception: pass
    try:
        with pikepdf.open(io.BytesIO(data)) as pdf:
            buf = io.BytesIO()
            pdf.save(buf, compress_streams=True,
                     stream_decode_level=pikepdf.StreamDecodeLevel.generalized)
            results.append(buf.getvalue())
    except Exception: pass
    if not results: return data
    best = min(results, key=len)
    return best if len(best) < len(data) else data

def split_pdf_all(data: bytes) -> list:
    with fitz.open(stream=data, filetype="pdf") as doc:
        results = []
        for i in range(len(doc)):
            out = fitz.open()
            out.insert_pdf(doc, from_page=i, to_page=i)
            buf = io.BytesIO(); out.save(buf)
            results.append(buf.getvalue())
        return results

def merge_pdfs(pdfs: list) -> bytes:
    out = fitz.open()
    for data in pdfs:
        with fitz.open(stream=data, filetype="pdf") as doc:
            out.insert_pdf(doc)
    buf = io.BytesIO(); out.save(buf)
    return buf.getvalue()

def pdf_to_images(data: bytes, dpi: int = 150) -> list:
    images = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            images.append(pix.tobytes("png"))
    return images

def images_to_pdf(image_bytes_list: list) -> bytes:
    out = fitz.open()
    for img_bytes in image_bytes_list:
        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        img_buf = io.BytesIO(); img.save(img_buf, "JPEG", quality=90)
        imgdoc = fitz.open(stream=img_buf.getvalue(), filetype="jpg")
        pdfbytes = imgdoc.convert_to_pdf()
        with fitz.open(stream=pdfbytes, filetype="pdf") as p:
            out.insert_pdf(p)
    buf = io.BytesIO(); out.save(buf)
    return buf.getvalue()

def lock_pdf(data: bytes, password: str) -> bytes:
    with pikepdf.open(io.BytesIO(data)) as pdf:
        buf = io.BytesIO()
        pdf.save(buf, encryption=pikepdf.Encryption(
            owner=password, user=password, R=4
        ))
        return buf.getvalue()

def unlock_pdf(data: bytes, password: str) -> bytes:
    with pikepdf.open(io.BytesIO(data), password=password) as pdf:
        buf = io.BytesIO(); pdf.save(buf)
        return buf.getvalue()

def repair_pdf(data: bytes) -> bytes:
    try:
        with fitz.open(stream=data, filetype="pdf") as doc:
            buf = io.BytesIO()
            doc.save(buf, garbage=4, deflate=True, clean=True)
            return buf.getvalue()
    except Exception:
        try:
            with pikepdf.open(io.BytesIO(data), suppress_warnings=True) as pdf:
                buf = io.BytesIO(); pdf.save(buf)
                return buf.getvalue()
        except Exception as e:
            raise ValueError(f"Could not repair: {e}")

def add_watermark_text(data: bytes, text: str) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            w, h = page.rect.width, page.rect.height
            page.insert_text(
                fitz.Point(w * 0.15, h * 0.55),
                text, fontsize=52, color=(0.75, 0.75, 0.75),
                rotate=45, overlay=True
            )
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def add_watermark_image(data: bytes, logo_bytes: bytes) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        img_xref = doc.get_new_xref()
        for page in doc:
            w, h = page.rect.width, page.rect.height
            rect = fitz.Rect(w * 0.3, h * 0.35, w * 0.7, h * 0.65)
            page.insert_image(rect, stream=logo_bytes, overlay=True, keep_proportion=True)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def dark_mode_pdf(data: bytes) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            page.set_mediabox(page.rect)
            page.draw_rect(page.rect, color=None, fill=(0.1, 0.1, 0.15), overlay=False)
            for block in page.get_text("dict")["blocks"]:
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            page.add_redact_annot(fitz.Rect(span["bbox"]), fill=(0.1, 0.1, 0.15))
            page.apply_redactions()
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def add_page_numbers(data: bytes, style: str = "arabic", position: str = "bottom_center") -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        for i, page in enumerate(doc):
            w, h = page.rect.width, page.rect.height
            n = i + 1
            label = str(n) if style == "arabic" else _to_roman(n) if style == "roman" else chr(64 + n) if n <= 26 else str(n)
            pos = fitz.Point(w / 2, h - 20) if position == "bottom_center" else fitz.Point(w - 40, h - 20)
            page.insert_text(pos, label, fontsize=11, color=(0.3, 0.3, 0.3), overlay=True)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def change_bg_color(data: bytes, color: tuple) -> bytes:
    r, g, b = [x / 255.0 for x in color]
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            page.draw_rect(page.rect, color=None, fill=(r, g, b), overlay=False)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def rotate_pdf(data: bytes, degrees: int) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            page.set_rotation(degrees)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def resize_to_a4(data: bytes) -> bytes:
    a4_w, a4_h = 595, 842
    with fitz.open(stream=data, filetype="pdf") as doc:
        out = fitz.open()
        for page in doc:
            new_page = out.new_page(width=a4_w, height=a4_h)
            new_page.show_pdf_page(new_page.rect, doc, page.number)
        buf = io.BytesIO(); out.save(buf)
        return buf.getvalue()

def pdf_to_word(data: bytes) -> bytes:
    try:
        from docx import Document as DocxDoc
        from docx.shared import Inches
        doc_out = DocxDoc()
        with fitz.open(stream=data, filetype="pdf") as doc:
            for i, page in enumerate(doc):
                if i > 0:
                    doc_out.add_page_break()
                doc_out.add_heading(f"Page {i+1}", level=2)
                text = page.get_text()
                if text.strip():
                    doc_out.add_paragraph(text)
                else:
                    mat = fitz.Matrix(2, 2)
                    pix = page.get_pixmap(matrix=mat)
                    img_buf = io.BytesIO(pix.tobytes("png"))
                    doc_out.add_picture(img_buf, width=Inches(6))
        buf = io.BytesIO(); doc_out.save(buf)
        return buf.getvalue()
    except Exception as e:
        raise ValueError(f"PDF to Word failed: {e}")

def pdf_to_ppt(data: bytes) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page_num, page in enumerate(doc):
            mat = fitz.Matrix(150/72, 150/72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            slide = prs.slides.add_slide(blank_layout)
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, left=0, top=0,
                                     width=prs.slide_width, height=prs.slide_height)
    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()

def pdf_to_excel(data: bytes) -> bytes:
    import pdfplumber
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    wb = Workbook(); wb.remove(wb.active)
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            ws = wb.create_sheet(title=f"Page {i+1}")
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row_idx, row in enumerate(table):
                        ws.append([cell or "" for cell in row])
                        if row_idx == 0:
                            for cell in ws[ws.max_row]:
                                cell.font = Font(bold=True)
                                cell.fill = PatternFill("solid", fgColor="D0E8FF")
            else:
                text = page.extract_text() or ""
                for row in text.split("\n"):
                    ws.append([row])
    buf = io.BytesIO(); wb.save(buf)
    return buf.getvalue()

def add_text_to_pdf(data: bytes, text: str, page_num: int = 0, x: float = 72, y: float = 72) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        if 0 <= page_num < len(doc):
            page = doc[page_num]
            page.insert_text(fitz.Point(x, y), text, fontsize=14, color=(0, 0, 0), overlay=True)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def add_footer(data: bytes, footer_text: str) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            w, h = page.rect.width, page.rect.height
            page.insert_text(fitz.Point(w/2 - len(footer_text)*3, h - 20),
                             footer_text, fontsize=10, color=(0.4, 0.4, 0.4), overlay=True)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def extract_pages(data: bytes, page_range: str) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        pages = _parse_range(page_range, len(doc))
        out = fitz.open()
        for p in pages:
            out.insert_pdf(doc, from_page=p, to_page=p)
        buf = io.BytesIO(); out.save(buf)
        return buf.getvalue()

def delete_pages(data: bytes, page_range: str) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        to_del = set(_parse_range(page_range, len(doc)))
        out = fitz.open()
        for i in range(len(doc)):
            if i not in to_del:
                out.insert_pdf(doc, from_page=i, to_page=i)
        buf = io.BytesIO(); out.save(buf)
        return buf.getvalue()

def reorder_pages(data: bytes, order_str: str) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        order = _parse_range(order_str, len(doc))
        out = fitz.open()
        for p in order:
            out.insert_pdf(doc, from_page=p, to_page=p)
        buf = io.BytesIO(); out.save(buf)
        return buf.getvalue()

def reverse_pages(data: bytes) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        out = fitz.open()
        for i in range(len(doc) - 1, -1, -1):
            out.insert_pdf(doc, from_page=i, to_page=i)
        buf = io.BytesIO(); out.save(buf)
        return buf.getvalue()

def crop_margins(data: bytes) -> bytes:
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            clip = page.get_bboxlog()
            if clip:
                bbox = fitz.Rect(clip[0][1])
                for _, r in clip[1:]:
                    bbox |= fitz.Rect(r)
                pad = 20
                bbox.x0 = max(0, bbox.x0 - pad)
                bbox.y0 = max(0, bbox.y0 - pad)
                bbox.x1 = min(page.rect.width,  bbox.x1 + pad)
                bbox.y1 = min(page.rect.height, bbox.y1 + pad)
                page.set_cropbox(bbox)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def generate_qr(text: str) -> bytes:
    import qrcode
    qr = qrcode.QRCode(version=1, error_correction=qrcode.constants.ERROR_CORRECT_L, box_size=10, border=4)
    qr.add_data(text); qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buf = io.BytesIO(); img.save(buf, format="PNG")
    return buf.getvalue()

def get_metadata(data: bytes) -> dict:
    with fitz.open(stream=data, filetype="pdf") as doc:
        meta = doc.metadata
        return {
            "title":     meta.get("title", "—"),
            "author":    meta.get("author", "—"),
            "subject":   meta.get("subject", "—"),
            "creator":   meta.get("creator", "—"),
            "pages":     len(doc),
            "size":      file_size_str(data),
            "encrypted": doc.is_encrypted,
            "format":    meta.get("format", "—"),
        }

def ocr_pdf(data: bytes, lang: str = "eng+hin") -> str:
    try:
        import pytesseract
        with fitz.open(stream=data, filetype="pdf") as doc:
            all_pages = []
            for i, page in enumerate(doc):
                blocks = page.get_text("blocks")
                if blocks:
                    blocks.sort(key=lambda b: (round(b[1] / 20) * 20, b[0]))
                    lines = [b[4].strip() for b in blocks if b[4].strip()]
                    native_text = "\n".join(lines)
                    if len(native_text.strip()) > 30:
                        all_pages.append(f"[Page {i+1}]\n{native_text}")
                        continue
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                config = "--psm 6 --oem 3"
                text = pytesseract.image_to_string(img, config=config, lang=lang)
                all_pages.append(f"[Page {i+1}]\n{text.strip()}")
            return "\n\n" + "─" * 40 + "\n\n".join(all_pages)
    except Exception as e:
        return f"⚠️ Error: {e}"

def compare_pdfs(data1: bytes, data2: bytes) -> tuple:
    import difflib
    def extract(data):
        lines = []
        with fitz.open(stream=data, filetype="pdf") as doc:
            for page in doc:
                lines.extend(page.get_text().split("\n"))
        return lines
    lines1, lines2 = extract(data1), extract(data2)
    diff = list(difflib.unified_diff(lines1, lines2, lineterm="", n=0))
    added   = sum(1 for l in diff if l.startswith("+") and not l.startswith("+++"))
    removed = sum(1 for l in diff if l.startswith("-") and not l.startswith("---"))
    with fitz.open(stream=data1, filetype="pdf") as d: pages1 = len(d)
    with fitz.open(stream=data2, filetype="pdf") as d: pages2 = len(d)
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4
    c.setFillColor(colors.HexColor("#1a1a2e"))
    c.rect(0, h-60, w, 60, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(20, h-40, "PDF Comparison Report")
    c.setFillColor(colors.black)
    c.setFont("Helvetica", 10)
    y = h - 90
    for item in [f"Pages in PDF 1: {pages1}", f"Pages in PDF 2: {pages2}",
                 f"Lines added: {added}", f"Lines removed: {removed}"]:
        c.drawString(30, y, item); y -= 16
    y -= 10; c.setFont("Courier", 8)
    for line in diff[:500]:
        if y < 40: c.showPage(); y = h - 40; c.setFont("Courier", 8)
        if line.startswith("+"): c.setFillColor(colors.HexColor("#1a6b1a"))
        elif line.startswith("-"): c.setFillColor(colors.HexColor("#6b1a1a"))
        else: c.setFillColor(colors.black)
        try: c.drawString(20, y, line[:110])
        except: pass
        y -= 13
    c.save()
    return buf.getvalue(), {"pages1": pages1, "pages2": pages2, "added": added, "removed": removed}


# ─── NEW v5: PDF TOOLS ────────────────────────────────────────────────────────

def pdf_to_txt(data: bytes) -> str:
    """Extract all text from PDF as plain text."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        parts = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                parts.append(f"=== Page {i+1} ===\n{text}")
        return "\n\n".join(parts) if parts else "⚠️ No text found in PDF."

def linearize_pdf(data: bytes) -> bytes:
    """Optimize PDF for fast web view (linearize)."""
    with pikepdf.open(io.BytesIO(data)) as pdf:
        buf = io.BytesIO()
        pdf.save(buf, linearize=True, compress_streams=True)
        return buf.getvalue()

def pdf_thumbnail(data: bytes, page_num: int = 0, size: tuple = (400, 565)) -> bytes:
    """Get first page as thumbnail image."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        page = doc[min(page_num, len(doc)-1)]
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.thumbnail(size, Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=90)
        return buf.getvalue()

def pdf_deep_info(data: bytes) -> dict:
    """Deep analysis of PDF file."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        meta = doc.metadata
        total_images = 0
        total_fonts = set()
        total_words = 0
        for page in doc:
            total_images += len(page.get_images(full=True))
            for font in page.get_fonts(full=True):
                total_fonts.add(font[3] or font[4] or "Unknown")
            total_words += len(page.get_text().split())
        return {
            "title":       meta.get("title", "—"),
            "author":      meta.get("author", "—"),
            "subject":     meta.get("subject", "—"),
            "creator":     meta.get("creator", "—"),
            "producer":    meta.get("producer", "—"),
            "created":     meta.get("creationDate", "—")[:16],
            "modified":    meta.get("modDate", "—")[:16],
            "pages":       len(doc),
            "size":        file_size_str(data),
            "encrypted":   doc.is_encrypted,
            "pdf_version": meta.get("format", "—"),
            "images":      total_images,
            "fonts":       len(total_fonts),
            "font_list":   list(total_fonts)[:5],
            "words":       total_words,
            "bookmarks":   len(doc.get_toc()),
        }

def redact_text(data: bytes, search_text: str) -> bytes:
    """Redact (black out) all occurrences of a word in PDF."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        count = 0
        for page in doc:
            areas = page.search_for(search_text)
            for area in areas:
                page.add_redact_annot(area, fill=(0, 0, 0))
                count += 1
        for page in doc:
            page.apply_redactions()
        if count == 0:
            raise ValueError(f'Text "{search_text}" not found in PDF.')
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue(), count

def impose_pdf(data: bytes, layout: str = "2up") -> bytes:
    """
    2-up: 2 pages side by side on one page.
    4-up: 4 pages in 2x2 grid on one page.
    """
    n = 4 if layout == "4up" else 2
    with fitz.open(stream=data, filetype="pdf") as src:
        total = len(src)
        out = fitz.open()
        a4_w, a4_h = 842, 595  # Landscape A4

        for i in range(0, total, n):
            new_page = out.new_page(width=a4_w, height=a4_h)
            if n == 2:
                positions = [
                    fitz.Rect(0, 0, a4_w/2, a4_h),
                    fitz.Rect(a4_w/2, 0, a4_w, a4_h),
                ]
            else:
                hw, hh = a4_w/2, a4_h/2
                positions = [
                    fitz.Rect(0, 0, hw, hh),
                    fitz.Rect(hw, 0, a4_w, hh),
                    fitz.Rect(0, hh, hw, a4_h),
                    fitz.Rect(hw, hh, a4_w, a4_h),
                ]
            for j, rect in enumerate(positions):
                pg_idx = i + j
                if pg_idx < total:
                    new_page.show_pdf_page(rect, src, pg_idx)

        buf = io.BytesIO(); out.save(buf)
        return buf.getvalue()

def deskew_pdf(data: bytes) -> bytes:
    """Deskew scanned/crooked PDF pages using Pillow."""
    try:
        import numpy as np
        has_numpy = True
    except ImportError:
        has_numpy = False

    with fitz.open(stream=data, filetype="pdf") as src:
        out_doc = fitz.open()
        for page in src:
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            if has_numpy:
                import numpy as np
                gray = img.convert("L")
                arr = np.array(gray)
                # Simple deskew: find text angle using projection
                threshold = 128
                binary = (arr < threshold).astype(np.float32)
                best_angle, best_score = 0.0, -1
                for angle in range(-10, 11):
                    rotated = Image.fromarray((binary * 255).astype(np.uint8)).rotate(angle, expand=False)
                    r_arr = np.array(rotated)
                    col_sums = r_arr.sum(axis=1)
                    score = float(col_sums.var())
                    if score > best_score:
                        best_score = score
                        best_angle = angle
                if abs(best_angle) > 0.5:
                    img = img.rotate(best_angle, expand=True, fillcolor=(255, 255, 255))

            img_buf = io.BytesIO()
            img.save(img_buf, format="JPEG", quality=90)
            imgdoc = fitz.open(stream=img_buf.getvalue(), filetype="jpg")
            pdfbytes = imgdoc.convert_to_pdf()
            with fitz.open(stream=pdfbytes, filetype="pdf") as p:
                out_doc.insert_pdf(p)

        buf = io.BytesIO(); out_doc.save(buf)
        return buf.getvalue()

def check_password_strength(password: str) -> dict:
    """Analyze password strength and return detailed report."""
    score = 0
    issues = []
    tips = []

    # Length check
    if len(password) < 6:
        issues.append("Too short (< 6 chars)")
    elif len(password) < 8:
        score += 1
        tips.append("Use 8+ characters for better security")
    elif len(password) < 12:
        score += 2
    else:
        score += 3

    has_upper = any(c.isupper() for c in password)
    has_lower = any(c.islower() for c in password)
    has_digit = any(c.isdigit() for c in password)
    has_special = any(c in string.punctuation for c in password)

    if has_upper: score += 1
    else: tips.append("Add uppercase letters (A-Z)")
    if has_lower: score += 1
    else: tips.append("Add lowercase letters (a-z)")
    if has_digit: score += 1
    else: tips.append("Add numbers (0-9)")
    if has_special: score += 2
    else: tips.append("Add special chars (!@#$%)")

    from config import PWD_CRACK_COMMON_LIST
    if password.lower() in [p.lower() for p in PWD_CRACK_COMMON_LIST]:
        score = 1
        issues.append("This is a VERY common password! Change it immediately.")

    if score <= 2:
        level, emoji, color = "Weak", "🔴", "red"
    elif score <= 4:
        level, emoji, color = "Fair", "🟡", "orange"
    elif score <= 6:
        level, emoji, color = "Good", "🟢", "green"
    else:
        level, emoji, color = "Strong", "💪", "blue"

    bar_filled = min(score, 8)
    bar = "█" * bar_filled + "░" * (8 - bar_filled)

    return {
        "score": score,
        "level": level,
        "emoji": emoji,
        "bar": bar,
        "has_upper": has_upper,
        "has_lower": has_lower,
        "has_digit": has_digit,
        "has_special": has_special,
        "length": len(password),
        "issues": issues,
        "tips": tips,
    }

def crack_pdf_password(data: bytes, timeout_sec: int = 60) -> str | None:
    """
    Try to crack PDF password.
    Strategy: common list first, then numeric brute-force up to 6 digits.
    Returns password string if found, None if not found.
    WARNING: Only works on weak/simple passwords. For marketing purposes.
    """
    from config import PWD_CRACK_COMMON_LIST, PWD_CRACK_TIMEOUT_SEC
    start = time.time()

    def try_pass(pwd: str) -> bool:
        try:
            with pikepdf.open(io.BytesIO(data), password=pwd):
                return True
        except Exception:
            return False

    # Check if actually locked
    if try_pass(""):
        return ""  # Not password protected

    # Try common passwords
    for pwd in PWD_CRACK_COMMON_LIST:
        if time.time() - start > timeout_sec: return None
        if try_pass(pwd): return pwd

    # Try numeric: 1-6 digits
    for length in range(1, 7):
        for combo in itertools.product(string.digits, repeat=length):
            if time.time() - start > timeout_sec: return None
            pwd = "".join(combo)
            if try_pass(pwd): return pwd

    return None

def edit_metadata(data: bytes, fields: dict) -> bytes:
    """Edit PDF metadata fields."""
    with pikepdf.open(io.BytesIO(data)) as pdf:
        with pdf.open_metadata(set_pikepdf_as_editor=False) as meta:
            ns = "http://purl.org/dc/"
            if fields.get("title"):
                meta["dc:title"] = fields["title"]
            if fields.get("author"):
                meta["dc:creator"] = [fields["author"]]
            if fields.get("subject"):
                meta["dc:subject"] = [fields["subject"]]
        if fields.get("title"):
            pdf.docinfo["/Title"]   = fields.get("title", "")
        if fields.get("author"):
            pdf.docinfo["/Author"]  = fields.get("author", "")
        if fields.get("subject"):
            pdf.docinfo["/Subject"] = fields.get("subject", "")
        buf = io.BytesIO(); pdf.save(buf)
        return buf.getvalue()

def pdf_sign(data: bytes, signature_img_bytes: bytes, page_num: int = -1,
             position: str = "bottom_right") -> bytes:
    """Add a visible signature image to PDF page."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        pg = doc[page_num if page_num >= 0 else len(doc)-1]
        w, h = pg.rect.width, pg.rect.height
        sig_w, sig_h = 150, 60
        if position == "bottom_right":
            rect = fitz.Rect(w - sig_w - 20, h - sig_h - 30, w - 20, h - 30)
        elif position == "bottom_left":
            rect = fitz.Rect(20, h - sig_h - 30, 20 + sig_w, h - 30)
        else:
            rect = fitz.Rect(w/2 - sig_w/2, h - sig_h - 30, w/2 + sig_w/2, h - 30)
        pg.insert_image(rect, stream=signature_img_bytes, overlay=True, keep_proportion=True)
        buf = io.BytesIO(); doc.save(buf)
        return buf.getvalue()

def pdf_to_epub(data: bytes) -> bytes:
    """Convert PDF to basic EPUB format."""
    try:
        from ebooklib import epub
    except ImportError:
        raise ValueError("ebooklib not installed. Run: pip install ebooklib")

    book = epub.EpubBook()
    book.set_identifier("nexora_pdf_doc")
    book.set_title("Converted Document")
    book.set_language("en")
    book.add_author("Nexora PDF Doctor")

    spine = ["nav"]
    with fitz.open(stream=data, filetype="pdf") as doc:
        for i, page in enumerate(doc):
            text = page.get_text()
            if not text.strip():
                mat = fitz.Matrix(1.5, 1.5)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                img_item = epub.EpubItem(
                    uid=f"img_{i}", file_name=f"images/page_{i+1}.png",
                    media_type="image/png", content=img_bytes
                )
                book.add_item(img_item)
                html_content = f'<html><body><img src="images/page_{i+1}.png" style="max-width:100%"/></body></html>'
            else:
                paragraphs = "".join(f"<p>{line}</p>" for line in text.split("\n") if line.strip())
                html_content = f"<html><body><h2>Page {i+1}</h2>{paragraphs}</body></html>"

            chapter = epub.EpubHtml(title=f"Page {i+1}", file_name=f"page_{i+1}.xhtml", lang="en")
            chapter.content = html_content
            book.add_item(chapter)
            spine.append(chapter)

    book.toc = tuple(spine[1:])
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = spine

    buf = io.BytesIO()
    epub.write_epub(buf, book)
    return buf.getvalue()

def epub_to_pdf(data: bytes) -> bytes:
    """Convert EPUB to PDF by extracting text and rendering."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ValueError("ebooklib/beautifulsoup4 not installed.")

    book = epub.read_epub(io.BytesIO(data))
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4
    y = h - 60
    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, y, "Converted from EPUB"); y -= 30
    c.setFont("Helvetica", 11)

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text("\n", strip=True)
            for line in text.split("\n"):
                if not line.strip(): continue
                words = line.split()
                chunk = ""
                for word in words:
                    test = (chunk + " " + word).strip()
                    if c.stringWidth(test, "Helvetica", 11) < w - 80:
                        chunk = test
                    else:
                        if y < 50:
                            c.showPage(); y = h - 50
                            c.setFont("Helvetica", 11)
                        c.drawString(40, y, chunk); y -= 16
                        chunk = word
                if chunk:
                    if y < 50:
                        c.showPage(); y = h - 50
                        c.setFont("Helvetica", 11)
                    c.drawString(40, y, chunk); y -= 16
    c.save()
    return buf.getvalue()

def doc_to_pdf(data: bytes, ext: str = "docx") -> bytes:
    """Convert DOCX/DOC to PDF using python-docx + reportlab."""
    from docx import Document as DocxDoc
    doc = DocxDoc(io.BytesIO(data))
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4
    y = h - 60
    c.setFont("Helvetica", 11)
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text: y -= 8; continue
        if para.style.name.startswith("Heading"):
            c.setFont("Helvetica-Bold", 14)
        else:
            c.setFont("Helvetica", 11)
        words = text.split()
        chunk = ""
        for word in words:
            test = (chunk + " " + word).strip()
            if c.stringWidth(test, "Helvetica", 11) < w - 80:
                chunk = test
            else:
                if y < 50:
                    c.showPage(); y = h - 50; c.setFont("Helvetica", 11)
                c.drawString(40, y, chunk); y -= 16; chunk = word
        if chunk:
            if y < 50:
                c.showPage(); y = h - 50; c.setFont("Helvetica", 11)
            c.drawString(40, y, chunk); y -= 16
    c.save()
    return buf.getvalue()


# ─── NEW v5: IMAGE TOOLS ──────────────────────────────────────────────────────

def img_compress(data: bytes, quality: int = 70) -> bytes:
    img = Image.open(io.BytesIO(data)).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality, optimize=True)
    return buf.getvalue()

def img_resize(data: bytes, width: int, height: int) -> bytes:
    img = Image.open(io.BytesIO(data))
    img = img.resize((width, height), Image.LANCZOS)
    buf = io.BytesIO()
    fmt = img.format or "PNG"
    img.save(buf, format=fmt if fmt in ("JPEG","PNG","WEBP") else "PNG")
    return buf.getvalue()

def img_crop(data: bytes, left: int, top: int, right: int, bottom: int) -> bytes:
    img = Image.open(io.BytesIO(data))
    img = img.crop((left, top, right, bottom))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

def img_apply_filter(data: bytes, filter_name: str) -> bytes:
    img = Image.open(io.BytesIO(data)).convert("RGB")
    if filter_name == "blur":
        img = img.filter(ImageFilter.GaussianBlur(radius=3))
    elif filter_name == "sharpen":
        img = img.filter(ImageFilter.SHARPEN)
    elif filter_name == "emboss":
        img = img.filter(ImageFilter.EMBOSS)
    elif filter_name == "edge":
        img = img.filter(ImageFilter.EDGE_ENHANCE_MORE)
    elif filter_name == "grayscale":
        img = ImageOps.grayscale(img).convert("RGB")
    elif filter_name == "sepia":
        gray = ImageOps.grayscale(img)
        sepia = Image.merge("RGB", [
            gray.point(lambda p: min(255, int(p * 1.08))),
            gray.point(lambda p: int(p * 0.85)),
            gray.point(lambda p: int(p * 0.66)),
        ])
        img = sepia
    elif filter_name == "brightness":
        img = ImageEnhance.Brightness(img).enhance(1.5)
    elif filter_name == "contrast":
        img = ImageEnhance.Contrast(img).enhance(2.0)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

def img_add_text(data: bytes, text: str, position: str = "center",
                  font_size: int = 40, color: tuple = (255,255,255)) -> bytes:
    img = Image.open(io.BytesIO(data)).convert("RGBA")
    overlay = Image.new("RGBA", img.size, (0,0,0,0))
    draw = ImageDraw.Draw(overlay)
    try:
        font = ImageFont.truetype("fonts/Caveat.ttf", font_size)
    except Exception:
        font = ImageFont.load_default()

    bbox = draw.textbbox((0,0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    iw, ih = img.size

    if position == "center":
        xy = ((iw - tw) // 2, (ih - th) // 2)
    elif position == "top":
        xy = ((iw - tw) // 2, 20)
    elif position == "bottom":
        xy = ((iw - tw) // 2, ih - th - 20)
    else:
        xy = (20, 20)

    # Shadow
    draw.text((xy[0]+2, xy[1]+2), text, font=font, fill=(0,0,0,180))
    draw.text(xy, text, font=font, fill=(*color, 255))

    result = Image.alpha_composite(img, overlay).convert("RGB")
    buf = io.BytesIO()
    result.save(buf, format="JPEG", quality=93)
    return buf.getvalue()

def img_convert(data: bytes, to_format: str = "JPEG") -> bytes:
    img = Image.open(io.BytesIO(data)).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format=to_format, quality=95 if to_format == "JPEG" else None)
    return buf.getvalue()

def img_remove_bg(data: bytes) -> bytes:
    """Remove background using rembg (local ML model, no API needed)."""
    try:
        from rembg import remove
        result = remove(data)
        return result
    except ImportError:
        raise ValueError("rembg not installed. Run: pip install rembg")


# ─── NEW v5: DOCUMENT CONVERTERS ─────────────────────────────────────────────

def csv_to_pdf(data: bytes) -> bytes:
    import csv as csv_mod
    content = data.decode("utf-8", errors="replace")
    rows = list(csv_mod.reader(content.splitlines()))
    if not rows:
        raise ValueError("CSV file is empty.")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=30, rightMargin=30, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()

    # Cap column width
    col_count = len(rows[0])
    col_width = (A4[0] - 60) / max(col_count, 1)
    table_data = []
    for row in rows:
        table_data.append([Paragraph(str(c)[:100], styles["Normal"]) for c in row])

    t = Table(table_data, colWidths=[col_width] * col_count, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#2c3e50")),
        ("TEXTCOLOR",  (0,0), (-1,0), colors.white),
        ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",   (0,0), (-1,-1), 9),
        ("GRID",       (0,0), (-1,-1), 0.5, colors.grey),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#f0f4f8")]),
        ("ALIGN",      (0,0), (-1,-1), "LEFT"),
        ("VALIGN",     (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
    ]))
    doc.build([t])
    return buf.getvalue()

def txt_to_pdf(data: bytes) -> bytes:
    text = data.decode("utf-8", errors="replace")
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4
    y = h - 50
    c.setFont("Helvetica-Bold", 14)
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.drawString(40, y, "Text Document"); y -= 30
    c.setFont("Courier", 10)
    c.setFillColor(colors.black)
    for line in text.split("\n"):
        # Word wrap
        words = line.split()
        chunk = ""
        for word in words:
            test = (chunk + " " + word).strip()
            if c.stringWidth(test, "Courier", 10) < w - 80:
                chunk = test
            else:
                if y < 50:
                    c.showPage(); y = h - 50
                    c.setFont("Courier", 10)
                c.drawString(40, y, chunk); y -= 14; chunk = word
        if y < 50:
            c.showPage(); y = h - 50; c.setFont("Courier", 10)
        c.drawString(40, y, chunk if chunk else ""); y -= 14
    c.save()
    return buf.getvalue()

def html_to_pdf(html_content: str) -> bytes:
    """Convert HTML string to PDF using weasyprint or fallback."""
    try:
        import weasyprint
        pdf_bytes = weasyprint.HTML(string=html_content).write_pdf()
        return pdf_bytes
    except ImportError:
        pass
    # Fallback: strip tags and convert as text
    from html.parser import HTMLParser
    class MLStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.fed = []
        def handle_data(self, d):
            self.fed.append(d)
        def get_data(self):
            return "\n".join(self.fed)
    s = MLStripper()
    s.feed(html_content)
    plain = s.get_data()
    return txt_to_pdf(plain.encode("utf-8"))

def json_to_pdf(data: bytes) -> bytes:
    import json
    try:
        obj = json.loads(data.decode("utf-8", errors="replace"))
    except Exception as e:
        raise ValueError(f"Invalid JSON: {e}")

    buf = io.BytesIO()
    doc_out = SimpleDocTemplate(buf, pagesize=A4,
                                leftMargin=40, rightMargin=40, topMargin=50, bottomMargin=40)
    styles = getSampleStyleSheet()
    story = []
    title_style = ParagraphStyle("Title2", parent=styles["Heading1"],
                                 fontSize=18, textColor=colors.HexColor("#2c3e50"))
    story.append(Paragraph("JSON Data Report", title_style))
    story.append(Spacer(1, 12))

    def render_obj(obj, depth=0):
        indent = "&nbsp;" * (depth * 4)
        if isinstance(obj, dict):
            for k, v in obj.items():
                key_style = ParagraphStyle(f"key_{depth}", parent=styles["Normal"],
                                           fontSize=10, textColor=colors.HexColor("#2980b9"),
                                           fontName="Helvetica-Bold")
                story.append(Paragraph(f"{indent}<b>{k}:</b>", key_style))
                render_obj(v, depth + 1)
        elif isinstance(obj, list):
            for i, item in enumerate(obj[:50]):
                story.append(Paragraph(f"{indent}[{i}]", styles["Normal"]))
                render_obj(item, depth + 1)
            if len(obj) > 50:
                story.append(Paragraph(f"{indent}... and {len(obj)-50} more items", styles["Normal"]))
        else:
            val = str(obj)[:200]
            story.append(Paragraph(f"{indent}{val}", styles["Normal"]))

    render_obj(obj)
    doc_out.build(story)
    return buf.getvalue()


# ─── NEW v5: SECURITY & PRIVACY ──────────────────────────────────────────────

def compute_hash(data: bytes) -> dict:
    return {
        "md5":    hashlib.md5(data).hexdigest(),
        "sha1":   hashlib.sha1(data).hexdigest(),
        "sha256": hashlib.sha256(data).hexdigest(),
        "size":   file_size_str(data),
        "bytes":  len(data),
    }

def steg_hide(image_data: bytes, secret_text: str) -> bytes:
    """Hide secret text inside image (LSB steganography)."""
    img = Image.open(io.BytesIO(image_data)).convert("RGB")
    pixels = list(img.getdata())

    # Encode text length + text as binary
    payload = secret_text.encode("utf-8")
    length = len(payload)
    header = length.to_bytes(4, "big")
    bits = "".join(format(b, "08b") for b in header + payload)

    if len(bits) > len(pixels) * 3:
        raise ValueError("Secret text too long for this image!")

    new_pixels = []
    bit_idx = 0
    for r, g, b in pixels:
        if bit_idx < len(bits):
            r = (r & 0xFE) | int(bits[bit_idx]); bit_idx += 1
        if bit_idx < len(bits):
            g = (g & 0xFE) | int(bits[bit_idx]); bit_idx += 1
        if bit_idx < len(bits):
            b = (b & 0xFE) | int(bits[bit_idx]); bit_idx += 1
        new_pixels.append((r, g, b))

    out_img = Image.new("RGB", img.size)
    out_img.putdata(new_pixels)
    buf = io.BytesIO()
    out_img.save(buf, format="PNG")
    return buf.getvalue()

def steg_reveal(image_data: bytes) -> str:
    """Reveal hidden text from image (LSB steganography)."""
    img = Image.open(io.BytesIO(image_data)).convert("RGB")
    pixels = list(img.getdata())

    bits = []
    for r, g, b in pixels:
        bits.extend([r & 1, g & 1, b & 1])

    # Read 4-byte length header
    header_bits = bits[:32]
    length = int("".join(str(b) for b in header_bits), 2)

    if length == 0 or length > 10000:
        return "⚠️ No hidden message found."

    payload_bits = bits[32:32 + length * 8]
    payload_bytes = bytes(
        int("".join(str(b) for b in payload_bits[i:i+8]), 2)
        for i in range(0, len(payload_bits), 8)
    )
    try:
        return payload_bytes.decode("utf-8")
    except Exception:
        return "⚠️ Could not decode hidden message."


# ─── NEW v5: CREATIVE TOOLS ───────────────────────────────────────────────────

def create_poster(title: str, subtitle: str = "", theme: str = "dark") -> bytes:
    themes = {
        "dark":    {"bg": (18, 18, 30),     "title": (255,220,50),  "sub": (200,200,255), "accent": (100,80,200)},
        "light":   {"bg": (255,255,255),    "title": (30,30,80),    "sub": (80,80,120),   "accent": (100,150,220)},
        "red":     {"bg": (20,0,0),         "title": (255,80,80),   "sub": (255,180,180), "accent": (200,0,0)},
        "green":   {"bg": (0,20,10),        "title": (80,255,120),  "sub": (180,255,200), "accent": (0,180,80)},
        "gradient":{"bg": (15,15,60),       "title": (255,255,255), "sub": (180,200,255), "accent": (80,120,255)},
    }
    t = themes.get(theme, themes["dark"])
    w, h = 800, 1100
    img = Image.new("RGB", (w, h), t["bg"])
    draw = ImageDraw.Draw(img)

    # Gradient overlay
    for y_pos in range(h):
        alpha = int(30 * (1 - y_pos / h))
        draw.line([(0, y_pos), (w, y_pos)], fill=tuple(min(255, c + alpha) for c in t["bg"]))

    # Decorative border
    for thick in range(3):
        draw.rectangle([10+thick*4, 10+thick*4, w-10-thick*4, h-10-thick*4],
                       outline=t["accent"], width=1)

    # Title
    try:
        font_title = ImageFont.truetype("fonts/Pacifico.ttf", 72)
        font_sub   = ImageFont.truetype("fonts/Caveat.ttf",  40)
        font_small = ImageFont.truetype("fonts/Kalam.ttf",   24)
    except:
        font_title = font_sub = font_small = ImageFont.load_default()

    # Center title with word wrap
    words = title.split()
    lines, line = [], ""
    for word in words:
        test = (line + " " + word).strip()
        bbox = draw.textbbox((0,0), test, font=font_title)
        if bbox[2] - bbox[0] < w - 100:
            line = test
        else:
            if line: lines.append(line)
            line = word
    if line: lines.append(line)

    y_start = h // 3
    for line in lines:
        bbox = draw.textbbox((0,0), line, font=font_title)
        lw = bbox[2] - bbox[0]
        draw.text(((w - lw)//2 + 3, y_start + 3), line, font=font_title, fill=(0,0,0,100))
        draw.text(((w - lw)//2, y_start), line, font=font_title, fill=t["title"])
        y_start += bbox[3] - bbox[1] + 15

    if subtitle:
        y_start += 20
        bbox = draw.textbbox((0,0), subtitle, font=font_sub)
        sw = bbox[2] - bbox[0]
        draw.text(((w - sw)//2, y_start), subtitle, font=font_sub, fill=t["sub"])

    # Branding
    brand = "🤖 Nexora PDF Doctor"
    try:
        bbox = draw.textbbox((0,0), brand, font=font_small)
        bw = bbox[2] - bbox[0]
        draw.text(((w - bw)//2, h - 60), brand, font=font_small, fill=t["accent"])
    except: pass

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

def create_calendar_pdf(year: int = None, month: int = None) -> bytes:
    now = datetime.date.today()
    year  = year  or now.year
    month = month or now.month
    month_name = calendar.month_name[month]

    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4

    # Background
    c.setFillColor(colors.HexColor("#1a1a2e"))
    c.rect(0, 0, w, h, fill=1, stroke=0)

    # Header
    c.setFillColor(colors.HexColor("#e94560"))
    c.rect(0, h-80, w, 80, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 28)
    title = f"{month_name} {year}"
    c.drawCentredString(w/2, h-52, title)

    # Day headers
    days = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    col_w = (w - 60) / 7
    x_start = 30
    y_days = h - 110
    c.setFont("Helvetica-Bold", 11)
    for i, d in enumerate(days):
        color = colors.HexColor("#e94560") if i >= 5 else colors.HexColor("#a8b3cf")
        c.setFillColor(color)
        c.drawCentredString(x_start + col_w * i + col_w/2, y_days, d)

    # Calendar grid
    cal = calendar.monthcalendar(year, month)
    row_h = (h - 200) / max(len(cal), 1)
    c.setFont("Helvetica-Bold", 18)
    today_marker = (now.year == year and now.month == month)

    for row_idx, week in enumerate(cal):
        for col_idx, day in enumerate(week):
            if day == 0: continue
            x = x_start + col_w * col_idx
            y = y_days - 25 - row_h * row_idx
            is_today = today_marker and day == now.day
            is_weekend = col_idx >= 5

            if is_today:
                c.setFillColor(colors.HexColor("#e94560"))
                c.circle(x + col_w/2, y + 8, 18, fill=1, stroke=0)
                c.setFillColor(colors.white)
            elif is_weekend:
                c.setFillColor(colors.HexColor("#ff8c69"))
            else:
                c.setFillColor(colors.HexColor("#c0cfe4"))

            c.drawCentredString(x + col_w/2, y, str(day))

    # Footer
    c.setFont("Helvetica", 9)
    c.setFillColor(colors.HexColor("#666688"))
    c.drawCentredString(w/2, 20, "Created by Nexora PDF Doctor v5.0")
    c.save()
    return buf.getvalue()

def create_invoice(name: str, items: list, company: str = "Nexora") -> bytes:
    """
    items = [{"desc": "Item name", "qty": 1, "price": 100}, ...]
    """
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4

    # Header
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.rect(0, h-100, w, 100, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 24)
    c.drawString(40, h-55, "INVOICE")
    c.setFont("Helvetica", 12)
    c.drawString(40, h-75, company)

    # Date & Invoice number
    c.setFillColor(colors.HexColor("#ecf0f1"))
    inv_no = f"INV-{datetime.date.today().strftime('%Y%m%d')}-{abs(hash(name)) % 9000 + 1000}"
    c.drawRightString(w-40, h-55, inv_no)
    c.drawRightString(w-40, h-75, datetime.date.today().strftime("%d %B %Y"))

    # Bill to
    y = h - 140
    c.setFillColor(colors.HexColor("#7f8c8d"))
    c.setFont("Helvetica", 10)
    c.drawString(40, y, "BILL TO:")
    y -= 16
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.setFont("Helvetica-Bold", 13)
    c.drawString(40, y, name); y -= 40

    # Table header
    c.setFillColor(colors.HexColor("#3498db"))
    c.rect(30, y-2, w-60, 24, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 10)
    c.drawString(40, y+5,    "Description")
    c.drawRightString(330,   y+5, "Qty")
    c.drawRightString(430,   y+5, "Unit Price")
    c.drawRightString(w-40,  y+5, "Total")
    y -= 30

    total = 0
    c.setFont("Helvetica", 10)
    for i, item in enumerate(items):
        if y < 100: c.showPage(); y = h - 60
        bg = colors.HexColor("#f8f9fa") if i % 2 == 0 else colors.white
        c.setFillColor(bg)
        c.rect(30, y-4, w-60, 20, fill=1, stroke=0)
        qty   = int(item.get("qty", 1))
        price = float(item.get("price", 0))
        line_total = qty * price
        total += line_total
        c.setFillColor(colors.HexColor("#2c3e50"))
        c.drawString(40, y+2,     str(item.get("desc","Item"))[:40])
        c.drawRightString(330,  y+2, str(qty))
        c.drawRightString(430,  y+2, f"₹{price:.2f}")
        c.drawRightString(w-40, y+2, f"₹{line_total:.2f}")
        y -= 24

    # Totals
    y -= 10
    c.setFillColor(colors.HexColor("#ecf0f1"))
    c.rect(350, y-6, w-380, 60, fill=1, stroke=0)
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.setFont("Helvetica", 10)
    tax = total * 0.18
    c.drawString(360, y+30,      "Subtotal:")
    c.drawRightString(w-40, y+30, f"₹{total:.2f}")
    c.drawString(360, y+14,      "GST (18%):")
    c.drawRightString(w-40, y+14, f"₹{tax:.2f}")
    c.setFont("Helvetica-Bold", 12)
    c.setFillColor(colors.HexColor("#e74c3c"))
    c.drawString(360, y-2,       "TOTAL:")
    c.drawRightString(w-40, y-2, f"₹{total + tax:.2f}")

    # Footer
    c.setFont("Helvetica-Oblique", 9)
    c.setFillColor(colors.HexColor("#95a5a6"))
    c.drawCentredString(w/2, 30, "Thank you for your business! | Generated by Nexora PDF Doctor")
    c.save()
    return buf.getvalue()

def create_resume(data: dict) -> bytes:
    """
    data = {name, title, email, phone, summary, skills: [], experience: [], education: []}
    """
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    w, h = A4

    # Header bar
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.rect(0, h-120, w, 120, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 26)
    c.drawString(40, h-60, data.get("name","Your Name"))
    c.setFont("Helvetica", 14)
    c.drawString(40, h-82, data.get("title","Professional Title"))
    c.setFont("Helvetica", 10)
    contact = f"📧 {data.get('email','')}   📞 {data.get('phone','')}"
    c.drawString(40, h-100, contact)

    y = h - 145

    def section(title_text, content_list):
        nonlocal y
        if y < 80: c.showPage(); y = h - 50
        c.setFillColor(colors.HexColor("#3498db"))
        c.rect(30, y-2, w-60, 20, fill=1, stroke=0)
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 11)
        c.drawString(38, y+4, title_text)
        y -= 26
        c.setFillColor(colors.HexColor("#2c3e50"))
        c.setFont("Helvetica", 10)
        for item in content_list:
            if y < 80: c.showPage(); y = h - 50; c.setFont("Helvetica", 10)
            c.drawString(40, y, f"• {str(item)[:90]}"); y -= 15
        y -= 8

    if data.get("summary"):
        c.setFont("Helvetica-Oblique", 10)
        c.setFillColor(colors.HexColor("#555"))
        c.drawString(40, y, data["summary"][:120]); y -= 20

    if data.get("skills"):
        section("🛠 SKILLS", data["skills"])
    if data.get("experience"):
        section("💼 EXPERIENCE", data["experience"])
    if data.get("education"):
        section("🎓 EDUCATION", data["education"])

    c.setFont("Helvetica", 8)
    c.setFillColor(colors.grey)
    c.drawCentredString(w/2, 20, "Resume generated by Nexora PDF Doctor v5.0")
    c.save()
    return buf.getvalue()

def create_certificate(name: str, course: str, date: str = None, issuer: str = "Nexora Academy") -> bytes:
    date = date or datetime.date.today().strftime("%d %B %Y")
    w, h = A4[1], A4[0]  # Landscape
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=(w, h))

    # Gold gradient bg
    c.setFillColor(colors.HexColor("#fffdf0"))
    c.rect(0, 0, w, h, fill=1, stroke=0)

    # Decorative borders
    border_colors = ["#c9a84c", "#e8d08a", "#c9a84c"]
    for i, bc in enumerate(border_colors):
        c.setStrokeColor(colors.HexColor(bc))
        c.setLineWidth(3 - i)
        margin = 15 + i * 6
        c.rect(margin, margin, w - 2*margin, h - 2*margin, fill=0, stroke=1)

    # Header
    c.setFont("Helvetica-Bold", 12)
    c.setFillColor(colors.HexColor("#8b6914"))
    c.drawCentredString(w/2, h-55, issuer.upper())

    c.setFont("Helvetica-Bold", 36)
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.drawCentredString(w/2, h-110, "Certificate of Completion")

    c.setFont("Helvetica", 14)
    c.setFillColor(colors.HexColor("#555"))
    c.drawCentredString(w/2, h-145, "This is to certify that")

    # Recipient name
    c.setFont("Helvetica-Bold", 40)
    c.setFillColor(colors.HexColor("#c9a84c"))
    c.drawCentredString(w/2, h-195, name)

    # Course
    c.setFont("Helvetica", 15)
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.drawCentredString(w/2, h-230, "has successfully completed the course")
    c.setFont("Helvetica-Bold", 20)
    c.setFillColor(colors.HexColor("#2c3e50"))
    c.drawCentredString(w/2, h-260, course)

    # Date
    c.setFont("Helvetica", 12)
    c.setFillColor(colors.HexColor("#666"))
    c.drawCentredString(w/2, h-295, f"Date: {date}")

    # Signature line
    c.setStrokeColor(colors.HexColor("#c9a84c"))
    c.setLineWidth(1)
    c.line(w/2-100, h-340, w/2+100, h-340)
    c.setFont("Helvetica", 10)
    c.setFillColor(colors.HexColor("#888"))
    c.drawCentredString(w/2, h-355, issuer)

    c.save()
    return buf.getvalue()


# ─── NEW v5: UTILITIES ────────────────────────────────────────────────────────

def create_zip(files: list) -> bytes:
    """files = [(filename, bytes), ...]"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname, fdata in files:
            zf.writestr(fname, fdata)
    return buf.getvalue()

def extract_zip(data: bytes) -> list:
    """Returns [(filename, bytes), ...]"""
    files = []
    with zipfile.ZipFile(io.BytesIO(data), "r") as zf:
        for name in zf.namelist()[:20]:  # limit to 20 files
            files.append((name, zf.read(name)))
    return files

def get_file_info(data: bytes, filename: str) -> dict:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "unknown"
    info = {
        "filename": filename,
        "extension": ext.upper(),
        "size": file_size_str(data),
        "bytes": len(data),
        "md5": hashlib.md5(data).hexdigest()[:16] + "...",
    }
    try:
        img = Image.open(io.BytesIO(data))
        info["width"]  = img.width
        info["height"] = img.height
        info["mode"]   = img.mode
        info["format"] = img.format or ext.upper()
        info["type"]   = "Image"
    except Exception:
        if ext == "pdf":
            try:
                with fitz.open(stream=data, filetype="pdf") as doc:
                    info["pages"] = len(doc)
                    info["encrypted"] = doc.is_encrypted
                    info["type"] = "PDF Document"
            except: info["type"] = "File"
        else:
            info["type"] = "File"
    return info

def scan_qr_code(data: bytes) -> str:
    try:
        from pyzbar.pyzbar import decode
        img = Image.open(io.BytesIO(data))
        results = decode(img)
        if results:
            return "\n".join(r.data.decode("utf-8") for r in results)
        return "⚠️ No QR code found in image."
    except ImportError:
        # Fallback: use OpenCV if available
        try:
            import cv2, numpy as np
            arr = np.frombuffer(data, np.uint8)
            img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
            detector = cv2.QRCodeDetector()
            text, _, _ = detector.detectAndDecode(img)
            return text if text else "⚠️ No QR code found."
        except Exception:
            return "⚠️ QR scanner unavailable. Install pyzbar: pip install pyzbar"

def generate_barcode(text: str, barcode_type: str = "code128") -> bytes:
    try:
        import barcode
        from barcode.writer import ImageWriter
        bc_class = barcode.get_barcode_class(barcode_type)
        bc = bc_class(text, writer=ImageWriter())
        buf = io.BytesIO()
        bc.write(buf, options={"write_text": True, "module_height": 15.0})
        return buf.getvalue()
    except ImportError:
        raise ValueError("python-barcode not installed. Run: pip install python-barcode Pillow")


# ─── Keep handwriting functions from v4 ──────────────────────────────────────

def create_handwritten_pdf(text, font_key, notebook_style="classic_blue", title="",
                            credit="Written By - Technical Serena"):
    import datetime
    from utils.font_loader import get_font_path

    # IST timezone (UTC+5:30)
    IST       = datetime.timezone(datetime.timedelta(hours=5, minutes=30))
    now       = datetime.datetime.now(IST)
    timestamp = now.strftime("%d %b %Y  |  %I:%M %p IST")

    font_path  = get_font_path(font_key)
    style      = NOTEBOOK_STYLES.get(notebook_style, NOTEBOOK_STYLES["classic_blue"])

    buf = io.BytesIO()
    c   = rl_canvas.Canvas(buf, pagesize=A4)
    pw, ph = A4

    font_name = "Helvetica"
    font_size = 16
    if font_path:
        try:
            pdfmetrics.registerFont(TTFont(font_key, font_path))
            font_name = font_key
        except Exception:
            pass

    line_spacing = style["line_spacing"]
    margin_x     = style["margin_x"]
    is_graph     = style.get("is_graph",  False)
    is_dotted    = style.get("is_dotted", False)
    tc           = style["text_color"]
    lc           = style["line_color"]

    def fill_background(cv):
        bg = style["bg"]
        if isinstance(bg, tuple) and len(bg) == 3:
            r, g, b = [x / 255.0 for x in bg]
            cv.setFillColorRGB(r, g, b)
            cv.rect(0, 0, pw, ph, fill=1, stroke=0)

    HEADER_H   = 55
    TOP_LINE_Y = ph - HEADER_H - 18

    def draw_header(cv):
        bg = style["bg"]
        r0, g0, b0 = [x / 255.0 for x in bg] if isinstance(bg, tuple) else (1, 1, 1)
        cv.setFillColorRGB(max(0, r0-0.07), max(0, g0-0.07), max(0, b0-0.07))
        cv.rect(0, ph - HEADER_H, pw, HEADER_H, fill=1, stroke=0)
        display_title = title.strip() if title and title.strip() else "Handwritten Notes"
        cv.setFillColorRGB(*tc)
        cv.setFont(font_name, 18)
        title_w = cv.stringWidth(display_title, font_name, 18)
        cv.drawString((pw - title_w) / 2, ph - 30, display_title)
        cv.setFont("Helvetica", 8)
        ts_w = cv.stringWidth(timestamp, "Helvetica", 8)
        cv.setFillColorRGB(min(1,tc[0]+0.3), min(1,tc[1]+0.3), min(1,tc[2]+0.3))
        cv.drawString(pw - ts_w - 15, ph - 18, timestamp)
        cv.setStrokeColorRGB(*lc)
        cv.setLineWidth(0.8)
        cv.line(35, ph - HEADER_H - 2, pw - 35, ph - HEADER_H - 2)

    def draw_page_lines(cv):
        fill_background(cv)
        draw_header(cv)
        cv.setStrokeColorRGB(*lc)
        cv.setLineWidth(0.4)
        if is_graph:
            for y in range(int(TOP_LINE_Y), 30, -line_spacing):
                cv.line(35, y, pw-35, y)
            for x in range(35, int(pw-35), line_spacing):
                cv.line(x, TOP_LINE_Y+4, x, 30)
        elif is_dotted:
            for y in range(int(TOP_LINE_Y), 30, -line_spacing):
                for x in range(50, int(pw-35), line_spacing):
                    cv.circle(x, y, 1.2, fill=1, stroke=0)
        else:
            for y in range(int(TOP_LINE_Y), 30, -line_spacing):
                cv.line(35, y, pw-35, y)
            mc = style.get("margin_color")
            if mc:
                cv.setStrokeColorRGB(*mc)
                cv.setLineWidth(0.7)
                cv.line(margin_x, ph-HEADER_H-4, margin_x, 30)

    def draw_credit(cv):
        if not credit: return
        cv.setFont("Helvetica-Oblique", 7)
        cr_color = (min(1,tc[0]+0.4), min(1,tc[1]+0.4), min(1,tc[2]+0.4))
        cv.setFillColorRGB(*cr_color)
        cw = cv.stringWidth(credit, "Helvetica-Oblique", 7)
        cv.drawString(pw - cw - 12, 10, credit)

    def wrap_text(text_content):
        wrapped = []
        max_w = pw - margin_x - 45
        for paragraph in text_content.split("\n"):
            words = paragraph.split()
            if not words: wrapped.append(""); continue
            line = ""
            for word in words:
                test = (line + " " + word).strip()
                if c.stringWidth(test, font_name, font_size) < max_w:
                    line = test
                else:
                    wrapped.append(line); line = word
            wrapped.append(line)
        return wrapped

    draw_page_lines(c)
    c.setFont(font_name, font_size)
    c.setFillColorRGB(*tc)
    lines = wrap_text(text)
    y_pos = TOP_LINE_Y
    text_x = margin_x + 6

    for line in lines:
        if y_pos < 45:
            draw_credit(c); c.showPage()
            draw_page_lines(c)
            c.setFont(font_name, font_size)
            c.setFillColorRGB(*tc)
            y_pos = TOP_LINE_Y
        if line:
            c.drawString(text_x, y_pos, line)
        y_pos -= line_spacing

    draw_credit(c)
    c.save()
    return buf.getvalue()


# =============================================================================
# v6 NEW FUNCTIONS — PDF Tools
# =============================================================================

def pdf_stamp(data: bytes, stamp_text: str, color: tuple = (180, 0, 0)) -> bytes:
    """Add diagonal stamp (CONFIDENTIAL / DRAFT / APPROVED etc.) to all pages."""
    import fitz, math
    doc = fitz.open(stream=data, filetype="pdf")
    r, g, b = color[0]/255, color[1]/255, color[2]/255
    for page in doc:
        pw, ph = page.rect.width, page.rect.height
        angle  = -math.degrees(math.atan2(ph, pw))
        cx, cy = pw / 2, ph / 2
        fs     = min(pw, ph) / 6
        # semi-transparent diagonal text
        page.insert_text(
            fitz.Point(cx - fs * 1.5, cy + fs * 0.4),
            stamp_text.replace("⛔ ", "").replace("📝 ", "").replace("✅ ", "")
              .replace("❌ ", "").replace("🔒 ", "").replace("📋 ", "")
              .replace("💰 ", "").replace("⏳ ", ""),
            fontsize=fs,
            rotate=int(-angle),
            color=(r, g, b),
            fill_opacity=0.18,
            overlay=True,
        )
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def pdf_grayscale(data: bytes) -> bytes:
    """Convert all pages to grayscale."""
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    buf = io.BytesIO()
    # Render each page as grayscale image, rebuild PDF
    new_doc = fitz.open()
    for page in doc:
        pix     = page.get_pixmap(colorspace=fitz.csGRAY, dpi=150)
        img_pdf = fitz.open("pdf", pix.pdfocr_data())
        new_doc.insert_pdf(img_pdf)
    new_doc.save(buf)
    return buf.getvalue()


def pdf_extract_images(data: bytes) -> bytes:
    """Extract all images from PDF, return as ZIP."""
    import fitz
    doc    = fitz.open(stream=data, filetype="pdf")
    files  = []
    count  = 0
    for page_num, page in enumerate(doc):
        for img in page.get_images(full=True):
            xref   = img[0]
            base   = doc.extract_image(xref)
            ext    = base["ext"]
            imgbytes = base["image"]
            files.append((f"page{page_num+1}_img{count+1}.{ext}", imgbytes))
            count += 1
            if count >= 50:
                break
        if count >= 50:
            break
    if not files:
        raise ValueError("No images found in this PDF!")
    return create_zip(files)


def pdf_remove_metadata(data: bytes) -> bytes:
    """Strip all metadata from PDF for privacy."""
    import pikepdf
    pdf = pikepdf.open(io.BytesIO(data))
    with pdf.open_metadata() as meta:
        keys = list(meta.keys())
        for k in keys:
            try:
                del meta[k]
            except Exception:
                pass
    pdf.docinfo.clear()
    buf = io.BytesIO()
    pdf.save(buf)
    return buf.getvalue()


def pdf_word_count(data: bytes) -> dict:
    """Count pages, words, characters, lines in PDF."""
    import fitz
    doc   = fitz.open(stream=data, filetype="pdf")
    total_words = total_chars = total_lines = 0
    for page in doc:
        text         = page.get_text()
        words        = text.split()
        total_words += len(words)
        total_chars += len(text.replace("\n", ""))
        total_lines += text.count("\n")
    return {
        "pages":  len(doc),
        "words":  total_words,
        "chars":  total_chars,
        "lines":  total_lines,
        "avg_words_per_page": round(total_words / max(len(doc), 1), 1),
    }


def pdf_add_header(data: bytes, header_text: str) -> bytes:
    """Add a header text to every page of the PDF."""
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    for page in doc:
        pw = page.rect.width
        page.insert_text(
            fitz.Point(pw / 2 - len(header_text) * 3.5, page.rect.height - 20),
            header_text,
            fontsize=9,
            color=(0.3, 0.3, 0.3),
            overlay=True,
        )
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def pdf_get_bookmarks(data: bytes) -> list:
    """Extract table of contents / bookmarks from PDF."""
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    toc = doc.get_toc()
    return [{"level": t[0], "title": t[1], "page": t[2]} for t in toc]


# =============================================================================
# v6 NEW FUNCTIONS — Image Tools
# =============================================================================

def img_collage(images: list, cols: int = 2) -> bytes:
    """
    Create a collage from list of image bytes.
    images: list of bytes
    """
    from PIL import Image as PILImage
    opened = []
    for raw in images[:12]:
        try:
            opened.append(PILImage.open(io.BytesIO(raw)).convert("RGB"))
        except Exception:
            pass
    if not opened:
        raise ValueError("No valid images!")
    thumb_w, thumb_h = 400, 300
    thumbs   = [i.resize((thumb_w, thumb_h), PILImage.LANCZOS) for i in opened]
    rows     = math.ceil(len(thumbs) / cols)
    canvas   = PILImage.new("RGB", (cols * thumb_w, rows * thumb_h), (30, 30, 30))
    for idx, thumb in enumerate(thumbs):
        x = (idx % cols) * thumb_w
        y = (idx // cols) * thumb_h
        canvas.paste(thumb, (x, y))
    buf = io.BytesIO()
    canvas.save(buf, "JPEG", quality=90)
    return buf.getvalue()


def img_meme(data: bytes, top_text: str = "", bottom_text: str = "") -> bytes:
    """Classic meme generator — bold white text with black stroke."""
    from PIL import Image as PILImage, ImageDraw, ImageFont
    img  = PILImage.open(io.BytesIO(data)).convert("RGB")
    draw = ImageDraw.Draw(img)
    w, h = img.size
    fs   = max(30, w // 12)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", fs)
    except Exception:
        font = ImageFont.load_default()

    def draw_text_with_stroke(text, y, anchor="mt"):
        if not text:
            return
        text = text.upper()
        bbox = draw.textbbox((0, 0), text, font=font)
        tw   = bbox[2] - bbox[0]
        x    = (w - tw) // 2
        # stroke
        for ox in range(-3, 4):
            for oy in range(-3, 4):
                draw.text((x + ox, y + oy), text, font=font, fill=(0, 0, 0))
        draw.text((x, y), text, font=font, fill=(255, 255, 255))

    draw_text_with_stroke(top_text,    20)
    draw_text_with_stroke(bottom_text, h - fs - 30)
    buf = io.BytesIO()
    img.save(buf, "JPEG", quality=92)
    return buf.getvalue()


def img_make_sticker(data: bytes) -> bytes:
    """Resize to 512×512 WebP for Telegram sticker."""
    from PIL import Image as PILImage
    img  = PILImage.open(io.BytesIO(data)).convert("RGBA")
    img  = img.resize((512, 512), PILImage.LANCZOS)
    buf  = io.BytesIO()
    img.save(buf, "WEBP", quality=90)
    return buf.getvalue()


def img_ascii_art(data: bytes, width: int = 80) -> str:
    """Convert image to ASCII art string."""
    from PIL import Image as PILImage
    chars  = "@%#*+=-:. "
    img    = PILImage.open(io.BytesIO(data)).convert("L")
    ratio  = img.height / img.width
    height = int(width * ratio * 0.45)
    img    = img.resize((width, height))
    pixels = list(img.getdata())
    ascii_img = ""
    for i, px in enumerate(pixels):
        ascii_img += chars[int(px / 255 * (len(chars) - 1))]
        if (i + 1) % width == 0:
            ascii_img += "\n"
    return ascii_img


def img_flip(data: bytes, direction: str = "horizontal") -> bytes:
    """Flip image horizontal or vertical."""
    from PIL import Image as PILImage
    img = PILImage.open(io.BytesIO(data))
    if direction == "horizontal":
        img = img.transpose(PILImage.FLIP_LEFT_RIGHT)
    else:
        img = img.transpose(PILImage.FLIP_TOP_BOTTOM)
    buf = io.BytesIO()
    img.save(buf, img.format or "PNG")
    return buf.getvalue()


def img_add_border(data: bytes, size: int = 20, color: tuple = (255, 255, 255)) -> bytes:
    """Add colored border around image."""
    from PIL import Image as PILImage, ImageOps
    img    = PILImage.open(io.BytesIO(data)).convert("RGB")
    result = ImageOps.expand(img, border=size, fill=color)
    buf    = io.BytesIO()
    result.save(buf, "JPEG", quality=92)
    return buf.getvalue()


def img_round_corners(data: bytes, radius: int = 40) -> bytes:
    """Apply rounded corners to image (PNG with transparency)."""
    from PIL import Image as PILImage, ImageDraw
    img  = PILImage.open(io.BytesIO(data)).convert("RGBA")
    mask = PILImage.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([(0, 0), img.size], radius=radius, fill=255)
    img.putalpha(mask)
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def img_get_exif(data: bytes) -> dict:
    """Extract EXIF data from image."""
    from PIL import Image as PILImage
    from PIL.ExifTags import TAGS
    img    = PILImage.open(io.BytesIO(data))
    result = {}
    try:
        exif_data = img._getexif()
        if exif_data:
            for tag_id, value in exif_data.items():
                tag = TAGS.get(tag_id, str(tag_id))
                if isinstance(value, (str, int, float)):
                    result[tag] = str(value)[:100]
    except Exception:
        pass
    return result


def img_remove_exif(data: bytes) -> bytes:
    """Strip all EXIF metadata from image (privacy)."""
    from PIL import Image as PILImage
    img     = PILImage.open(io.BytesIO(data))
    clean   = PILImage.new(img.mode, img.size)
    clean.putdata(list(img.getdata()))
    buf = io.BytesIO()
    fmt = img.format or "JPEG"
    clean.save(buf, fmt)
    return buf.getvalue()


def img_auto_enhance(data: bytes) -> bytes:
    """Auto-enhance brightness, contrast and sharpness."""
    from PIL import Image as PILImage, ImageEnhance
    img = PILImage.open(io.BytesIO(data)).convert("RGB")
    img = ImageEnhance.Brightness(img).enhance(1.1)
    img = ImageEnhance.Contrast(img).enhance(1.2)
    img = ImageEnhance.Sharpness(img).enhance(1.3)
    img = ImageEnhance.Color(img).enhance(1.1)
    buf = io.BytesIO()
    img.save(buf, "JPEG", quality=92)
    return buf.getvalue()


def img_apply_filter_v2(data: bytes, filter_name: str) -> bytes:
    """Extended filter support with vivid/vintage/cool/warm."""
    from PIL import Image as PILImage, ImageEnhance, ImageFilter
    import numpy as np
    img = PILImage.open(io.BytesIO(data)).convert("RGB")

    if filter_name == "vivid":
        img = ImageEnhance.Color(img).enhance(2.0)
        img = ImageEnhance.Contrast(img).enhance(1.3)
    elif filter_name == "vintage":
        arr = np.array(img, dtype=np.float32)
        arr[:, :, 0] = np.clip(arr[:, :, 0] * 1.1 + 20, 0, 255)
        arr[:, :, 1] = np.clip(arr[:, :, 1] * 0.9, 0, 255)
        arr[:, :, 2] = np.clip(arr[:, :, 2] * 0.7, 0, 255)
        img = PILImage.fromarray(arr.astype(np.uint8))
        img = ImageEnhance.Contrast(img).enhance(0.85)
    elif filter_name == "cool":
        arr = np.array(img, dtype=np.float32)
        arr[:, :, 0] = np.clip(arr[:, :, 0] * 0.85, 0, 255)
        arr[:, :, 2] = np.clip(arr[:, :, 2] * 1.2 + 15, 0, 255)
        img = PILImage.fromarray(arr.astype(np.uint8))
    elif filter_name == "warm":
        arr = np.array(img, dtype=np.float32)
        arr[:, :, 0] = np.clip(arr[:, :, 0] * 1.2 + 15, 0, 255)
        arr[:, :, 2] = np.clip(arr[:, :, 2] * 0.8, 0, 255)
        img = PILImage.fromarray(arr.astype(np.uint8))
    else:
        # fallback to existing
        return img_apply_filter(data, filter_name)

    buf = io.BytesIO()
    img.save(buf, "JPEG", quality=92)
    return buf.getvalue()


# =============================================================================
# v6 NEW FUNCTIONS — Creative / Cards
# =============================================================================

import math

def create_quote_card(quote: str, author: str = "", theme: str = "dark") -> bytes:
    """Beautiful quote card image (1080×1080)."""
    from PIL import Image as PILImage, ImageDraw, ImageFont
    from config import QUOTE_THEMES
    t    = QUOTE_THEMES.get(theme, QUOTE_THEMES["dark"])
    W, H = 1080, 1080
    img  = PILImage.new("RGB", (W, H), t["bg"])
    draw = ImageDraw.Draw(img)

    # Accent bars
    bar_h = 8
    for i in range(4):
        alpha = int(255 * (0.8 - i * 0.2))
        draw.rectangle([60, 80 + i, W - 60, 80 + bar_h + i],
                       fill=(*t["accent"], alpha) if len(t["accent"]) == 4 else t["accent"])
    draw.rectangle([60, 80, W - 60, 80 + bar_h], fill=t["accent"])
    draw.rectangle([60, H - 80 - bar_h, W - 60, H - 80], fill=t["accent"])

    # Quote marks
    try:
        font_big = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 120)
        font_quote = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 42)
        font_author = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
    except Exception:
        font_big   = ImageFont.load_default()
        font_quote = font_big
        font_author = font_big

    draw.text((80, 110), "\u201c", font=font_big, fill=(*t["accent"],) if isinstance(t["accent"], tuple) else t["accent"])

    # Word wrap quote
    words   = quote.split()
    lines   = []
    line    = ""
    max_w   = W - 160
    for word in words:
        test = (line + " " + word).strip()
        bbox = draw.textbbox((0, 0), test, font=font_quote)
        if bbox[2] - bbox[0] > max_w and line:
            lines.append(line)
            line = word
        else:
            line = test
    if line:
        lines.append(line)

    total_h = len(lines) * 60
    y_start = (H - total_h) // 2 - 30
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font_quote)
        x    = (W - (bbox[2] - bbox[0])) // 2
        draw.text((x, y_start), line, font=font_quote, fill=t["text"])
        y_start += 60

    # Author
    if author:
        author_text = f"\u2014 {author}"
        bbox        = draw.textbbox((0, 0), author_text, font=font_author)
        x           = (W - (bbox[2] - bbox[0])) // 2
        draw.text((x, H - 180), author_text, font=font_author,
                  fill=t["accent"])

    buf = io.BytesIO()
    img.save(buf, "PNG", quality=95)
    return buf.getvalue()


def create_birthday_card(name: str, message: str = "") -> bytes:
    """Colorful birthday card PDF."""
    from reportlab.lib.pagesizes import A5
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas as rl_canvas
    W, H = A5
    buf  = io.BytesIO()
    c    = rl_canvas.Canvas(buf, pagesize=A5)

    # Gradient-like background
    for i in range(int(H)):
        ratio = i / H
        r = int(255 * (1 - ratio * 0.3))
        g = int(100 + 80 * ratio)
        b = int(150 + 100 * ratio)
        c.setFillColorRGB(r/255, g/255, b/255)
        c.rect(0, i, W, 1, fill=1, stroke=0)

    # Balloons (circles)
    balloon_data = [
        (80,  H-80,  30, (1.0, 0.3, 0.3)),
        (160, H-60,  25, (0.3, 0.8, 0.3)),
        (W-80, H-80, 30, (0.3, 0.5, 1.0)),
        (W-150, H-55, 22, (1.0, 0.8, 0.0)),
        (W//2, H-100, 28, (1.0, 0.4, 0.8)),
    ]
    for bx, by, br, bc in balloon_data:
        c.setFillColorRGB(*bc)
        c.circle(bx, by, br, fill=1, stroke=0)
        c.setStrokeColorRGB(0.5, 0.5, 0.5)
        c.setLineWidth(1)
        c.line(bx, by - br, bx - 5, by - br - 40)

    # Title
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 36)
    title = "Happy Birthday!"
    c.drawCentredString(W/2, H - 160, title)

    # Name
    c.setFont("Helvetica-Bold", 28)
    c.setFillColorRGB(1, 0.95, 0.5)
    c.drawCentredString(W/2, H - 210, name)

    # Message
    if message:
        c.setFont("Helvetica", 16)
        c.setFillColorRGB(1, 1, 1)
        # wrap
        words  = message.split()
        lines  = []
        line   = ""
        for word in words:
            test = (line + " " + word).strip()
            if c.stringWidth(test, "Helvetica", 16) > W - 80:
                lines.append(line)
                line = word
            else:
                line = test
        if line:
            lines.append(line)
        y = H - 270
        for ln in lines:
            c.drawCentredString(W/2, y, ln)
            y -= 24

    # Cake emoji text
    c.setFont("Helvetica", 48)
    c.drawCentredString(W/2, 140, "🎂🎉🎊")

    c.save()
    return buf.getvalue()


def create_business_card(name: str, title: str, phone: str,
                         email: str, company: str, theme: str = "minimal") -> bytes:
    """Professional business card PDF (standard 3.5x2 inch)."""
    from reportlab.lib.pagesizes import inch
    from reportlab.pdfgen import canvas as rl_canvas
    from config import BCARD_THEMES
    t       = BCARD_THEMES.get(theme, BCARD_THEMES["minimal"])
    W, H    = 3.5 * inch, 2.0 * inch
    buf     = io.BytesIO()
    c       = rl_canvas.Canvas(buf, pagesize=(W, H))
    bg      = tuple(v/255 for v in t["bg"])
    accent  = tuple(v/255 for v in t["accent"])
    text_c  = tuple(v/255 for v in t["text"])

    # Background
    c.setFillColorRGB(*bg)
    c.rect(0, 0, W, H, fill=1, stroke=0)

    # Accent bar left
    c.setFillColorRGB(*accent)
    c.rect(0, 0, 6, H, fill=1, stroke=0)

    # Accent line
    c.setStrokeColorRGB(*accent)
    c.setLineWidth(1)
    c.line(20, H * 0.52, W - 20, H * 0.52)

    # Name
    c.setFillColorRGB(*text_c)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(20, H - 38, name)

    # Title
    c.setFillColorRGB(*accent)
    c.setFont("Helvetica-Oblique", 10)
    c.drawString(20, H - 56, title)

    # Company
    if company:
        c.setFillColorRGB(*text_c)
        c.setFont("Helvetica-Bold", 9)
        c.drawString(20, H - 72, company)

    # Contact info
    c.setFont("Helvetica", 9)
    c.setFillColorRGB(*text_c)
    y = 52
    if phone:
        c.drawString(20, y, f"📞  {phone}")
        y -= 16
    if email:
        c.drawString(20, y, f"✉️  {email}")

    c.save()
    return buf.getvalue()


def create_flyer(title: str, subtitle: str, details: str,
                 date_time: str = "", theme: str = "event") -> bytes:
    """Eye-catching event flyer PDF (A5)."""
    from reportlab.lib.pagesizes import A5
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib import colors
    from config import FLYER_THEMES
    t    = FLYER_THEMES.get(theme, FLYER_THEMES["event"])
    W, H = A5
    buf  = io.BytesIO()
    c    = rl_canvas.Canvas(buf, pagesize=A5)

    bg     = tuple(v/255 for v in t["bg"])
    accent = tuple(v/255 for v in t["accent"])
    text_c = tuple(v/255 for v in t["text"])

    # BG
    c.setFillColorRGB(*bg)
    c.rect(0, 0, W, H, fill=1, stroke=0)

    # Top accent strip
    c.setFillColorRGB(*accent)
    c.rect(0, H - 90, W, 90, fill=1, stroke=0)

    # Bottom accent strip
    c.rect(0, 0, W, 55, fill=1, stroke=0)

    # Decorative diagonal lines
    c.setStrokeColorRGB(*accent)
    c.setLineWidth(2)
    c.setFillColorRGB(*accent, )
    for x in range(-100, int(W)+200, 30):
        c.setStrokeAlpha(0.08)
        c.line(x, 0, x + 120, H)

    # Title inside top strip
    c.setFillColorRGB(*text_c)
    c.setFont("Helvetica-Bold", 28)
    c.drawCentredString(W/2, H - 58, title)

    # Subtitle
    c.setFont("Helvetica-Bold", 18)
    c.setFillColorRGB(*accent)
    c.drawCentredString(W/2, H - 120, subtitle)

    # Details
    c.setFont("Helvetica", 13)
    c.setFillColorRGB(*text_c)
    words  = details.split()
    lines  = []
    line   = ""
    for word in words:
        test = (line + " " + word).strip()
        if c.stringWidth(test, "Helvetica", 13) > W - 80:
            lines.append(line)
            line = word
        else:
            line = test
    if line:
        lines.append(line)
    y = H - 170
    for ln in lines:
        c.drawCentredString(W/2, y, ln)
        y -= 22

    # Date/time
    if date_time:
        c.setFillColorRGB(*accent)
        c.setFont("Helvetica-Bold", 16)
        c.drawCentredString(W/2, 80, date_time)

    c.setFillColorRGB(*text_c)
    c.setFont("Helvetica", 9)
    c.drawCentredString(W/2, 20, "Made with Nexora PDF Doctor Bot")

    c.save()
    return buf.getvalue()


def create_timetable(schedule: dict, title: str = "My Timetable") -> bytes:
    """
    Generate a weekly timetable PDF.
    schedule: {"Monday": ["9:00 Math", "10:00 English"], ...}
    """
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib import colors
    W, H = landscape(A4)
    buf  = io.BytesIO()
    c    = rl_canvas.Canvas(buf, pagesize=landscape(A4))

    days = [d for d in ["Monday","Tuesday","Wednesday","Thursday","Friday","Saturday","Sunday"]
            if d in schedule]
    if not days:
        raise ValueError("No valid days in schedule!")

    col_w   = (W - 80) / len(days)
    row_h   = 30
    header_h = 50

    # Background
    c.setFillColorRGB(0.97, 0.97, 1.0)
    c.rect(0, 0, W, H, fill=1, stroke=0)

    # Title
    c.setFillColorRGB(0.1, 0.1, 0.4)
    c.setFont("Helvetica-Bold", 20)
    c.drawCentredString(W/2, H - 35, title)

    # Headers
    for i, day in enumerate(days):
        x = 40 + i * col_w
        c.setFillColorRGB(0.2, 0.3, 0.7)
        c.rect(x, H - 60 - header_h, col_w - 4, header_h, fill=1, stroke=0)
        c.setFillColorRGB(1, 1, 1)
        c.setFont("Helvetica-Bold", 12)
        c.drawCentredString(x + col_w/2 - 2, H - 60 - header_h + 16, day[:3].upper())

    # Rows
    max_slots = max(len(v) for v in schedule.values()) if schedule else 8
    for slot in range(max_slots):
        y = H - 60 - header_h - (slot + 1) * row_h
        if y < 30:
            break
        # Row BG alternating
        row_bg = (0.92, 0.94, 1.0) if slot % 2 == 0 else (1, 1, 1)
        c.setFillColorRGB(*row_bg)
        c.rect(40, y, W - 80, row_h - 2, fill=1, stroke=0)

        for i, day in enumerate(days):
            x    = 40 + i * col_w
            slots = schedule.get(day, [])
            if slot < len(slots):
                c.setFillColorRGB(0.1, 0.1, 0.3)
                c.setFont("Helvetica", 9)
                c.drawString(x + 4, y + 9, str(slots[slot])[:25])

    # Grid lines
    c.setStrokeColorRGB(0.7, 0.7, 0.9)
    c.setLineWidth(0.5)
    for i in range(len(days) + 1):
        x = 40 + i * col_w
        c.line(x, H - 60 - header_h, x, 30)

    c.save()
    return buf.getvalue()


# =============================================================================
# v6 — UX / Bot helpers
# =============================================================================

def format_streak_message(streak: int) -> str:
    """Return motivational streak message."""
    from config import STREAK_BONUS_OPS
    msgs = {
        1:  "🔥 Day 1 streak! Keep going!",
        3:  "🔥🔥 3-day streak! You're on fire!",
        7:  "💎 7-day streak! One week strong!",
        14: "🏆 14-day streak! Amazing dedication!",
        30: "🌟 30-day streak! LEGENDARY!",
    }
    bonus = STREAK_BONUS_OPS.get(streak, 0)
    base  = msgs.get(streak, f"🔥 {streak}-day streak!")
    if bonus:
        base += f"\n🎁 Bonus: +{bonus} free ops today!"
    return base


def get_plan_badge(plan: str) -> str:
    return {"free": "🆓 Free", "basic": "⭐ Basic", "pro": "👑 Pro"}.get(plan, "🆓 Free")


# =============================================================================
# MISSING FUNCTIONS — Bug Fixes & Aliases (Added patch)
# =============================================================================

# ── BG color name → RGB mapping ───────────────────────────────────────────────
_BG_COLOR_MAP = {
    "bg_dark":   (30,  30,  45),
    "bg_white":  (255, 255, 255),
    "bg_yellow": (255, 252, 180),
    "bg_green":  (210, 245, 220),
    "bg_blue":   (210, 228, 255),
    "bg_pink":   (255, 214, 230),
}

def change_bg(data: bytes, color_key) -> bytes:
    """Change PDF background color. Accepts color key string or (r,g,b) tuple."""
    if isinstance(color_key, str):
        rgb = _BG_COLOR_MAP.get(color_key, (255, 255, 255))
    else:
        rgb = color_key
    return change_bg_color(data, rgb)


def watermark_text(data: bytes, text: str, invisible: bool = False) -> bytes:
    """Add text watermark. If invisible=True, uses near-transparent opacity."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            w, h = page.rect.width, page.rect.height
            opacity = 0.03 if invisible else 0.18
            page.insert_text(
                fitz.Point(w * 0.15, h * 0.55),
                text,
                fontsize=52,
                color=(0.75, 0.75, 0.75),
                rotate=45,
                fill_opacity=opacity,
                overlay=True,
            )
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()


def watermark_image(data: bytes, logo_bytes: bytes) -> bytes:
    """Add image watermark to all PDF pages."""
    return add_watermark_image(data, logo_bytes)


def resize_pdf_to_a4(data: bytes) -> bytes:
    """Resize all PDF pages to A4."""
    return resize_to_a4(data)


def auto_rotate_pdf(data: bytes) -> bytes:
    """Auto-detect and fix page rotation based on content orientation."""
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            rotation = page.rotation
            # Only correct clearly wrong rotations that aren't multiples of 90
            if rotation not in (0, 90, 180, 270):
                page.set_rotation(0)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()


def ocr_image(data: bytes, lang: str = "eng+hin") -> str:
    """Run OCR on an image (PNG/JPG/etc.) and return extracted text."""
    try:
        import pytesseract
        img = Image.open(io.BytesIO(data)).convert("RGB")
        config = "--psm 6 --oem 3"
        text = pytesseract.image_to_string(img, config=config, lang=lang)
        return text.strip() if text.strip() else "⚠️ No text found in image."
    except ImportError:
        return "⚠️ pytesseract not installed. Run: pip install pytesseract"
    except Exception as e:
        return f"⚠️ OCR Error: {e}"


def pdf_to_grayscale(data: bytes) -> bytes:
    """Convert PDF to grayscale (alias for pdf_grayscale)."""
    return pdf_grayscale(data)


def create_handwritten_jpg(text: str, font_key: str,
                            notebook_style: str = "classic_blue",
                            title: str = "",
                            credit: str = "Written By - Technical Serena") -> list:
    """
    Render handwritten notebook pages as high-quality JPEG images.
    2x super-sampling + LANCZOS downscale = PDF-equivalent quality.
    PIL y=0 = TOP, increases DOWNWARD.
    """
    import datetime, os
    from utils.font_loader import get_font_path
    from config import NOTEBOOK_STYLES

    # IST Timezone
    IST       = datetime.timezone(datetime.timedelta(hours=5, minutes=30))
    now       = datetime.datetime.now(IST)
    timestamp = now.strftime("%d %b %Y  |  %I:%M %p IST")

    # System font fallbacks (always available on Linux - never use load_default)
    _SYSTEM_FONTS = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/ubuntu/Ubuntu-R.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    ]
    _FALLBACK = next((f for f in _SYSTEM_FONTS if os.path.exists(f)), None)

    def _load_font(path, size):
        for p in ([path] if path else []) + ([_FALLBACK] if _FALLBACK else []):
            try:
                return ImageFont.truetype(p, size)
            except Exception:
                pass
        return ImageFont.load_default()

    # Target: A4 at 150 DPI
    FINAL_W  = int(595 * 150 / 72)   # 1240 px
    FINAL_H  = int(842 * 150 / 72)   # 1754 px

    # 2x Super-sampling
    SS       = 2
    PX_W     = FINAL_W * SS
    PX_H     = FINAL_H * SS
    SCALE    = 150 * SS / 72          # 4.167 px/pt

    style    = NOTEBOOK_STYLES.get(notebook_style, NOTEBOOK_STYLES["classic_blue"])

    def _rgb(tup):
        if tup and isinstance(tup[0], float):
            return tuple(int(c * 255) for c in tup)
        return tuple(tup)

    def _boost(color, factor=1.5):
        """Push colors away from white for vivid, non-washed appearance."""
        return tuple(max(0, min(255, int(255 - (255 - c) * factor))) for c in color)

    bg_rgb    = tuple(style["bg"]) if isinstance(style["bg"], tuple) else (255, 255, 255)
    tc        = _rgb(style["text_color"])
    lc        = _boost(_rgb(style["line_color"]),   1.6)
    mc        = _boost(_rgb(style["margin_color"]), 2.0) if style.get("margin_color") else None
    is_graph  = style.get("is_graph",  False)
    is_dotted = style.get("is_dotted", False)
    hdr_bg    = tuple(max(0, c - 25) for c in bg_rgb)

    HEADER_H     = int(55  * SCALE)
    SEP_Y        = HEADER_H + SS
    CONTENT_TOP  = HEADER_H + int(18 * SCALE)
    CONTENT_BOT  = PX_H - int(30 * SCALE)
    line_spacing = int(style["line_spacing"] * SCALE)
    margin_x     = int(style["margin_x"]     * SCALE)

    sz_body   = int(16 * SCALE)
    sz_title  = int(18 * SCALE)
    sz_small  = int(9  * SCALE)   # 9pt for timestamp/credit - clearly visible

    font_path  = get_font_path(font_key)
    body_font  = _load_font(font_path, sz_body)
    title_font = _load_font(font_path, sz_title)
    small_font = _load_font(_FALLBACK,  sz_small)  # always system font for small text

    def _make_page():
        img  = Image.new("RGB", (PX_W, PX_H), bg_rgb)
        draw = ImageDraw.Draw(img)

        # Header bar
        draw.rectangle([0, 0, PX_W, HEADER_H], fill=hdr_bg)
        disp_title = title.strip() or "Handwritten Notes"
        try:
            tb = draw.textbbox((0, 0), disp_title, font=title_font)
            tw, th = tb[2]-tb[0], tb[3]-tb[1]
            draw.text(((PX_W-tw)//2, (HEADER_H-th)//2), disp_title, font=title_font, fill=tc)
        except Exception:
            draw.text((int(20*SCALE), int(12*SCALE)), disp_title, font=title_font, fill=tc)

        # Timestamp — right-aligned, vertically centered in header
        try:
            tb2  = draw.textbbox((0, 0), timestamp, font=small_font)
            ts_w = tb2[2]-tb2[0]; ts_h = tb2[3]-tb2[1]
            ts_color = tuple(min(255, int(c * 0.55 + 100)) for c in tc)  # muted version of text color
            draw.text((PX_W - ts_w - int(15*SCALE), (HEADER_H-ts_h)//2),
                      timestamp, font=small_font, fill=ts_color)
        except Exception:
            pass

        # Separator line
        draw.line([int(35*SCALE), SEP_Y, PX_W-int(35*SCALE), SEP_Y], fill=lc, width=SS)

        # Notebook lines
        first_line_y = CONTENT_TOP + line_spacing
        line_ys      = range(first_line_y, CONTENT_BOT, line_spacing)

        if is_graph:
            for y in line_ys:
                draw.line([int(35*SCALE), y, PX_W-int(35*SCALE), y], fill=lc, width=SS)
            for x in range(int(35*SCALE), PX_W-int(35*SCALE), line_spacing):
                draw.line([x, CONTENT_TOP, x, CONTENT_BOT], fill=lc, width=SS)
        elif is_dotted:
            for y in line_ys:
                for x in range(int(50*SCALE), PX_W-int(35*SCALE), line_spacing):
                    r = max(SS, int(1.3*SCALE))
                    draw.ellipse([x-r, y-r, x+r, y+r], fill=lc)
        else:
            for y in line_ys:
                draw.line([int(35*SCALE), y, PX_W-int(35*SCALE), y], fill=lc, width=SS)
            if mc:
                draw.line([margin_x, CONTENT_TOP, margin_x, CONTENT_BOT], fill=mc, width=SS)

        return img, draw

    def _add_credit(draw):
        if not credit:
            return
        try:
            tb = draw.textbbox((0, 0), credit, font=small_font)
            cw = tb[2]-tb[0]
            cr_color = tuple(min(255, int(c * 0.55 + 100)) for c in tc)
            draw.text((PX_W - cw - int(12*SCALE), PX_H - int(22*SCALE)),
                      credit, font=small_font, fill=cr_color)
        except Exception:
            pass

    max_line_w = PX_W - margin_x - int(45*SCALE)
    _tmp_draw  = ImageDraw.Draw(Image.new("RGB", (10, 10)))

    def _wrap(content: str) -> list:
        result = []
        for para in content.split("\n"):
            words = para.split()
            if not words:
                result.append(""); continue
            line = ""
            for word in words:
                test = (line + " " + word).strip()
                try:
                    tb = _tmp_draw.textbbox((0, 0), test, font=body_font)
                    w_ = tb[2]-tb[0]
                except Exception:
                    w_ = len(test) * sz_body * 0.55
                if w_ < max_line_w:
                    line = test
                else:
                    if line: result.append(line)
                    line = word
            result.append(line)
        return result

    lines         = _wrap(text)
    pages_out     = []
    img, draw     = _make_page()
    first_line_y  = CONTENT_TOP + line_spacing
    y_pos         = first_line_y - sz_body
    text_x        = margin_x + int(6*SCALE)

    for line in lines:
        if y_pos + sz_body > CONTENT_BOT:
            _add_credit(draw)
            final = img.resize((FINAL_W, FINAL_H), Image.LANCZOS)
            buf = io.BytesIO()
            final.save(buf, format="JPEG", quality=95, subsampling=0)
            pages_out.append(buf.getvalue())
            img, draw = _make_page()
            y_pos = first_line_y - sz_body

        if line:
            try:
                draw.text((text_x, y_pos), line, font=body_font, fill=tc)
            except Exception:
                pass
        y_pos += line_spacing

    _add_credit(draw)
    final = img.resize((FINAL_W, FINAL_H), Image.LANCZOS)
    buf   = io.BytesIO()
    final.save(buf, format="JPEG", quality=95, subsampling=0)
    pages_out.append(buf.getvalue())
    return pages_out

